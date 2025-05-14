import argparse
import json
import sys
import os
from pptx import Presentation

def extract_transcript(pptx_file, cue_token="transition", step_id="extract_transcript"):
    """
    Extracts speaker notes from a PPTX file and segments them based on a cue token.

    Args:
        pptx_file (str): Path to the PPTX file.
        cue_token (str): Token used to split segments in speaker notes.
        step_id (str): Identifier for the extraction step.

    Returns:
        dict: A dictionary containing extraction results.
    """
    manifest = {
        "step_id": step_id,
        "status": "failure",
        "pptx": pptx_file,
        "slide_count": 0,
        "slides": []
    }

    try:
        prs = Presentation(pptx_file)
        manifest["slide_count"] = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            slide_data = {"index": i, "segments": []}
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame is None or not notes_slide.notes_text_frame.text:
                raise ValueError(f"Slide {i} has no speaker notes.")

            notes_text = notes_slide.notes_text_frame.text
            segments = notes_text.split(cue_token)

            for j, segment_text in enumerate(segments):
                segment_text = segment_text.strip()
                if j < len(segments) - 1:
                    slide_data["segments"].append({"kind": "text", "text": segment_text})
                    slide_data["segments"].append({"kind": "cue", "cue": cue_token})
                else:
                    if segment_text:
                        slide_data["segments"].append({"kind": "text", "text": segment_text})

            manifest["slides"].append(slide_data)

        manifest["status"] = "success"

    except Exception as e:
        manifest["status"] = "failure"
        manifest["error"] = str(e)
        print(json.dumps(manifest, indent=2))
        sys.exit(1)

    return manifest

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract speaker notes from PPTX.")
    parser.add_argument("--pptx", required=True, help="Path to the PPTX file.")
    parser.add_argument("--out", help="Path to output manifest JSON file.")
    parser.add_argument("--cue", default="transition", help="Token to split segments.")
    parser.add_argument("--step_id", default="extract_transcript", help="Identifier for the step.")

    args = parser.parse_args()

    manifest_data = extract_transcript(args.pptx, args.cue, args.step_id)

    print(json.dumps(manifest_data, indent=2))

    if args.out:
        output_dir = os.path.dirname(args.out)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        with open(args.out, "w") as f:
            json.dump(manifest_data, f, indent=2)

    if manifest_data["status"] == "failure":
        sys.exit(1)
