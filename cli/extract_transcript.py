import argparse
import json
import sys
import os
import re
from pptx import Presentation

def _process_segment_text(text_segment):
    """Helper function to process text segments for ellipses."""
    segment_text = text_segment.strip()
    segment_data = {"kind": "text"}
    leading_ellipse = False
    trailing_ellipse = False

    if segment_text.startswith("..."):
        leading_ellipse = True
        segment_text = segment_text[3:].lstrip()
    elif segment_text.startswith("…"): # Unicode ellipsis
        leading_ellipse = True
        segment_text = segment_text[1:].lstrip()

    if segment_text.endswith("..."):
        trailing_ellipse = True
        segment_text = segment_text[:-3].rstrip()
    elif segment_text.endswith("…"): # Unicode ellipsis
        trailing_ellipse = True
        segment_text = segment_text[:-1].rstrip()

    segment_data["text"] = segment_text

    if leading_ellipse and trailing_ellipse:
        segment_data["continue"] = ["start", "end"]
    elif leading_ellipse:
        segment_data["continue"] = "start"
    elif trailing_ellipse:
        segment_data["continue"] = "end"
    
    return segment_data, segment_text # Return segment_text to check if it's empty

def extract_transcript(pptx_file, cue_token="[transition]", step_id="extract_transcript"):
    """
    Extracts speaker notes from a PPTX file and segments them based on a cue token.

    Args:
        pptx_file (str): Path to the PPTX file.
        cue_token (str): Token used to split segments in speaker notes. Case-insensitive.
                         The exact cue_token string (preserving its case) will be used in the output.
        step_id (str): Identifier for the extraction step.

    Returns:
        dict: A dictionary containing extraction results.
    """
    manifest = {
        "step_id": step_id,
        "status": "failure",
        "pptx": pptx_file,
        "slide_count": 0,
        "slides": []
    }

    try:
        prs = Presentation(pptx_file)
        manifest["slide_count"] = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            slide_data = {"index": i, "segments": []}
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame is None or not notes_slide.notes_text_frame.text:
                # If there are no notes, we can either raise an error or return an empty segments list.
                # Current behavior is to raise ValueError. For robustness, perhaps allow empty notes.
                # For now, keeping existing behavior:
                # raise ValueError(f"Slide {i} has no speaker notes.")
                # Or, to allow slides with no notes (e.g. title slides):
                manifest["slides"].append(slide_data) # Add slide with empty segments
                continue


            notes_text = notes_slide.notes_text_frame.text
            
            current_pos = 0
            # Find all non-overlapping occurrences of the cue_token, case-insensitively
            # re.escape is important if cue_token contains special regex characters
            for match in re.finditer(re.escape(cue_token), notes_text, flags=re.IGNORECASE):
                match_start, match_end = match.span()
                
                # Add text segment before the cue
                text_before_cue = notes_text[current_pos:match_start]
                if text_before_cue.strip(): # Process only if there's non-whitespace text
                    processed_segment_data, processed_text = _process_segment_text(text_before_cue)
                    if processed_text: # Add segment only if it's not empty after stripping ellipses and whitespace
                        slide_data["segments"].append(processed_segment_data)
                
                # Add the cue segment, using the originally defined cue_token
                slide_data["segments"].append({"kind": "cue", "cue": cue_token})
                current_pos = match_end
            
            # Add any remaining text after the last cue
            text_after_last_cue = notes_text[current_pos:]
            if text_after_last_cue.strip(): # Process only if there's non-whitespace text
                processed_segment_data, processed_text = _process_segment_text(text_after_last_cue)
                if processed_text: # Add segment only if it's not empty
                    slide_data["segments"].append(processed_segment_data)

            # If no cues were found and the original notes_text was not empty, 
            # the whole notes_text is a single segment.
            # This check ensures we don't add an empty text segment if notes_text was only whitespace.
            if not slide_data["segments"] and notes_text.strip():
                processed_segment_data, processed_text = _process_segment_text(notes_text)
                if processed_text:
                    slide_data["segments"].append(processed_segment_data)
                
            manifest["slides"].append(slide_data)

        manifest["status"] = "success"

    except Exception as e:
        manifest["status"] = "failure"
        manifest["error"] = str(e)
        # print(json.dumps(manifest, indent=2)) # Avoid printing here if used as a library
        # sys.exit(1) # Avoid sys.exit if used as a library
        raise # Re-raise the exception so the caller can handle it or it terminates script

    return manifest

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract speaker notes from PPTX.")
    parser.add_argument("--pptx", required=True, help="Path to the PPTX file.")
    parser.add_argument("--out", help="Path to output manifest JSON file.")
    parser.add_argument("--cue", default="[transition]", help="Token to split segments.")
    parser.add_argument("--step_id", default="extract_transcript", help="Identifier for the step.")

    args = parser.parse_args()

    try:
        manifest_data = extract_transcript(args.pptx, args.cue, args.step_id)
        print(json.dumps(manifest_data, indent=2))

        if args.out:
            output_dir = os.path.dirname(args.out)
            if output_dir and not os.path.exists(output_dir): # Create dir only if it has a name
                os.makedirs(output_dir, exist_ok=True)
            with open(args.out, "w") as f:
                json.dump(manifest_data, f, indent=2)

        if manifest_data["status"] == "failure":
            sys.exit(1)
    except Exception as e:
        # Construct a failure manifest if extract_transcript raises an error
        # before it can construct its own.
        error_manifest = {
            "step_id": args.step_id,
            "status": "failure",
            "pptx": args.pptx,
            "error": str(e)
        }
        print(json.dumps(error_manifest, indent=2))
        sys.exit(1)
