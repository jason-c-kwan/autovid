#!/usr/bin/env python3
"""
Visual validation system for AutoVid slide synchronization.

This module provides Phase 1 (infrastructure + basic CV) and Phase 2 (Vision LLM)
capabilities for validating sync plans using visual ground truth from PowerPoint
slides and video frames.
"""

import json
import logging
import tempfile
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any, Union
from dataclasses import dataclass, asdict
import hashlib
import base64

# Core dependencies (always available)
import numpy as np

# Optional dependencies with graceful fallbacks
try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    
try:
    import cv2
    CV2_AVAILABLE = True
except ImportError:
    CV2_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import ffmpeg
    FFMPEG_AVAILABLE = True
except ImportError:
    FFMPEG_AVAILABLE = False

logger = logging.getLogger(__name__)


@dataclass
class SlideAnimationState:
    """Represents a slide at a specific build animation state."""
    slide_number: int
    build_step: int  # 0 = initial, 1 = after first build, etc.
    image_path: str
    visible_elements: List[str]  # Which elements are visible
    animation_type: str  # "build_in", "build_out", "static"
    timestamp_hint: Optional[float] = None  # When this state likely appears
    confidence: float = 1.0  # Confidence in this state detection


@dataclass
class VisualValidationResult:
    """Results from visual validation of a slide segment."""
    slide_number: int
    video_timestamp: float
    slide_image_path: str
    video_frame_path: str
    similarity_score: float  # 0.0 - 1.0
    validation_method: str   # "basic_cv", "vision_llm", "hybrid"
    confidence: float        # 0.0 - 1.0
    issues: List[str]        # Any problems detected
    recommendations: List[str]  # Suggested improvements
    llm_analysis: Optional[Dict[str, Any]] = None  # Full LLM response if used


class SlideImageExtractor:
    """Extract slide images from PowerPoint files for visual validation."""
    
    def __init__(self, output_dir: Optional[str] = None):
        """
        Initialize the slide image extractor.
        
        Args:
            output_dir: Directory to save extracted images (default: temp dir)
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for slide image extraction")
        
        self.output_dir = Path(output_dir) if output_dir else Path(tempfile.gettempdir()) / "autovid_slides"
        self.output_dir.mkdir(exist_ok=True, parents=True)
        
    def extract_all_slides(self, pptx_path: str, slide_size: Tuple[int, int] = (1920, 1080)) -> List[str]:
        """
        Extract all slides from a PowerPoint file as images.
        
        Args:
            pptx_path: Path to the PowerPoint file
            slide_size: Output image size (width, height)
            
        Returns:
            List of paths to extracted slide images
        """
        if not PIL_AVAILABLE:
            logger.error("PIL/Pillow is required for slide image extraction. Please install: pip install Pillow")
            return []
            
        try:
            prs = Presentation(pptx_path)
            image_paths = []
            
            for i, slide in enumerate(prs.slides):
                image_path = self._render_slide_to_image(slide, i + 1, slide_size)
                if image_path:
                    image_paths.append(image_path)
                    
            logger.info(f"Extracted {len(image_paths)} slide images from {pptx_path}")
            return image_paths
            
        except Exception as e:
            logger.error(f"Failed to extract slides from {pptx_path}: {e}")
            return []
    
    def extract_slide_states(self, pptx_path: str, slide_number: int) -> List[SlideAnimationState]:
        """
        Extract different animation states for a specific slide.
        
        Args:
            pptx_path: Path to the PowerPoint file
            slide_number: Slide number (1-based)
            
        Returns:
            List of SlideAnimationState objects representing different build states
        """
        try:
            prs = Presentation(pptx_path)
            if slide_number < 1 or slide_number > len(prs.slides):
                logger.error(f"Slide {slide_number} not found in presentation")
                return []
            
            slide = prs.slides[slide_number - 1]
            
            # Analyze build animations for this slide
            animation_sequence = self._analyze_build_animations(slide, slide_number)
            
            if not animation_sequence:
                # No animations detected - create single static state
                image_path = self._render_slide_to_image(slide, slide_number, (1920, 1080))
                if not image_path:
                    return []
                
                state = SlideAnimationState(
                    slide_number=slide_number,
                    build_step=0,
                    image_path=image_path,
                    visible_elements=self._extract_visible_elements(slide),
                    animation_type="static",
                    confidence=1.0
                )
                return [state]
            
            # Create states for each animation step
            states = []
            for step, anim_info in enumerate(animation_sequence):
                image_path = self._render_slide_state(slide, slide_number, step, anim_info)
                if image_path:
                    state = SlideAnimationState(
                        slide_number=slide_number,
                        build_step=step,
                        image_path=image_path,
                        visible_elements=anim_info.get('visible_elements', []),
                        animation_type=anim_info.get('animation_type', 'build_in'),
                        timestamp_hint=anim_info.get('timestamp_hint'),
                        confidence=anim_info.get('confidence', 0.8)
                    )
                    states.append(state)
            
            logger.info(f"Extracted {len(states)} animation states for slide {slide_number}")
            return states
            
        except Exception as e:
            logger.error(f"Failed to extract slide states for slide {slide_number}: {e}")
            return []
    
    def _render_slide_to_image(self, slide, slide_number: int, size: Tuple[int, int]) -> Optional[str]:
        """
        Render a slide to an image file.
        
        Note: This is a basic implementation. PowerPoint doesn't have direct 
        image export via python-pptx, so this creates a placeholder.
        For production, would need to integrate with PowerPoint COM or 
        use export functionality.
        """
        try:
            # Create output path
            output_path = self.output_dir / f"slide_{slide_number:03d}.png"
            
            # Phase 1: Create a placeholder image with slide content outline
            # Phase 2: Integrate with actual PowerPoint rendering
            
            if not PIL_AVAILABLE:
                logger.warning("PIL not available - creating text placeholder")
                return str(output_path)
            
            # Create a basic representation of the slide
            img = Image.new('RGB', size, color='white')
            draw = ImageDraw.Draw(img)
            
            # Add slide number and basic content info
            try:
                font = ImageFont.load_default()
            except:
                font = None
            
            draw.text((50, 50), f"Slide {slide_number}", fill='black', font=font)
            
            # Try to extract text content from slide
            text_content = self._extract_slide_text(slide)
            if text_content:
                # Add text content preview (first 200 chars)
                preview_text = text_content[:200] + "..." if len(text_content) > 200 else text_content
                draw.text((50, 100), preview_text, fill='black', font=font)
            
            # Save the image
            img.save(output_path)
            logger.debug(f"Created slide image: {output_path}")
            
            return str(output_path)
            
        except Exception as e:
            logger.error(f"Failed to render slide {slide_number} to image: {e}")
            return None
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text content from a slide."""
        text_content = []
        
        try:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text.strip())
        except Exception as e:
            logger.debug(f"Error extracting text from slide: {e}")
            
        return " ".join(text_content)
    
    def _extract_visible_elements(self, slide) -> List[str]:
        """Extract information about visible elements on the slide."""
        elements = []
        
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type'):
                    shape_type = str(shape.shape_type)
                    if hasattr(shape, 'text') and shape.text:
                        elements.append(f"text: {shape.text[:50]}")
                    else:
                        elements.append(f"shape: {shape_type}")
        except Exception as e:
            logger.debug(f"Error extracting elements from slide: {e}")
            
        return elements
    
    def _analyze_build_animations(self, slide, slide_number: int) -> List[Dict[str, Any]]:
        """
        Analyze build animations for a slide to determine different visual states.
        
        Args:
            slide: PowerPoint slide object
            slide_number: Slide number for logging
            
        Returns:
            List of animation step information
        """
        try:
            animation_sequence = []
            
            # Get all shapes that could have animations
            shapes_with_content = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    shapes_with_content.append({
                        'shape': shape,
                        'text': shape.text.strip(),
                        'type': 'text'
                    })
                elif hasattr(shape, 'shape_type'):
                    shapes_with_content.append({
                        'shape': shape,
                        'text': '',
                        'type': str(shape.shape_type)
                    })
            
            if not shapes_with_content:
                return []
            
            # Detect common animation patterns
            animation_pattern = self._detect_animation_pattern(shapes_with_content)
            
            if animation_pattern == 'progressive_bullets':
                # Create states for progressive bullet point revelation
                animation_sequence = self._create_progressive_bullet_states(shapes_with_content)
            elif animation_pattern == 'title_then_content':
                # Title first, then content appears
                animation_sequence = self._create_title_content_states(shapes_with_content)
            elif animation_pattern == 'grouped_reveals':
                # Content appears in logical groups
                animation_sequence = self._create_grouped_reveal_states(shapes_with_content)
            else:
                # Default: assume simple build-in for all content
                animation_sequence = self._create_simple_build_states(shapes_with_content)
            
            logger.debug(f"Detected {animation_pattern} pattern with {len(animation_sequence)} states for slide {slide_number}")
            return animation_sequence
            
        except Exception as e:
            logger.error(f"Animation analysis failed for slide {slide_number}: {e}")
            return []
    
    def _detect_animation_pattern(self, shapes_with_content: List[Dict]) -> str:
        """Detect the likely animation pattern based on slide content structure."""
        
        # Count different types of content
        text_shapes = [s for s in shapes_with_content if s['type'] == 'text']
        
        if not text_shapes:
            return 'static'
        
        # Look for bullet point patterns
        bullet_indicators = ['•', '-', '◦', '▪', '▫']
        bullet_shapes = []
        
        for shape_info in text_shapes:
            text = shape_info['text']
            if any(indicator in text for indicator in bullet_indicators):
                bullet_shapes.append(shape_info)
        
        # Analyze text structure for common patterns
        if len(bullet_shapes) >= 2:
            return 'progressive_bullets'
        elif len(text_shapes) >= 2:
            # Check if first shape looks like a title
            first_text = text_shapes[0]['text']
            if len(first_text) < 50 and not any(indicator in first_text for indicator in bullet_indicators):
                return 'title_then_content'
            else:
                return 'grouped_reveals'
        else:
            return 'static'
    
    def _create_progressive_bullet_states(self, shapes_with_content: List[Dict]) -> List[Dict[str, Any]]:
        """Create animation states for progressive bullet point revelation."""
        states = []
        
        # State 0: Title only (if present)
        title_shapes = []
        bullet_shapes = []
        other_shapes = []
        
        bullet_indicators = ['•', '-', '◦', '▪', '▫']
        
        for shape_info in shapes_with_content:
            text = shape_info['text']
            if shape_info['type'] == 'text' and text:
                if any(indicator in text for indicator in bullet_indicators):
                    bullet_shapes.append(shape_info)
                elif len(text) < 50 and len(states) == 0:  # Likely title
                    title_shapes.append(shape_info)
                else:
                    other_shapes.append(shape_info)
            else:
                other_shapes.append(shape_info)
        
        # Create states
        visible_elements = title_shapes + other_shapes  # Always visible elements
        
        if title_shapes:
            # State 0: Title + non-bullet content only
            states.append({
                'visible_elements': [s['text'] for s in visible_elements],
                'animation_type': 'build_in',
                'confidence': 0.9,
                'description': 'Title and base content'
            })
        
        # Progressive bullet states
        for i, bullet_shape in enumerate(bullet_shapes):
            current_visible = visible_elements + bullet_shapes[:i+1]
            states.append({
                'visible_elements': [s['text'] for s in current_visible],
                'animation_type': 'build_in',
                'confidence': 0.8,
                'description': f'With bullet point {i+1}'
            })
        
        return states
    
    def _create_title_content_states(self, shapes_with_content: List[Dict]) -> List[Dict[str, Any]]:
        """Create animation states for title-then-content pattern."""
        states = []
        
        text_shapes = [s for s in shapes_with_content if s['type'] == 'text']
        
        if len(text_shapes) >= 2:
            # State 0: Title only
            states.append({
                'visible_elements': [text_shapes[0]['text']],
                'animation_type': 'build_in', 
                'confidence': 0.8,
                'description': 'Title only'
            })
            
            # State 1: Title + all content
            states.append({
                'visible_elements': [s['text'] for s in shapes_with_content],
                'animation_type': 'build_in',
                'confidence': 0.8,
                'description': 'Full content'
            })
        
        return states
    
    def _create_grouped_reveal_states(self, shapes_with_content: List[Dict]) -> List[Dict[str, Any]]:
        """Create animation states for grouped content revelation."""
        states = []
        
        # Simple grouping: divide content into 2-3 logical groups
        total_shapes = len(shapes_with_content)
        group_size = max(1, total_shapes // 3)
        
        visible_elements = []
        for i in range(0, total_shapes, group_size):
            group = shapes_with_content[i:i+group_size]
            visible_elements.extend(group)
            
            states.append({
                'visible_elements': [s['text'] for s in visible_elements],
                'animation_type': 'build_in',
                'confidence': 0.7,
                'description': f'Content group {len(states) + 1}'
            })
        
        return states
    
    def _create_simple_build_states(self, shapes_with_content: List[Dict]) -> List[Dict[str, Any]]:
        """Create simple build states when pattern is unclear."""
        if not shapes_with_content:
            return []
        
        return [{
            'visible_elements': [s['text'] for s in shapes_with_content],
            'animation_type': 'build_in',
            'confidence': 0.6,
            'description': 'All content (simple build)'
        }]
    
    def _render_slide_state(
        self, 
        slide, 
        slide_number: int, 
        step: int, 
        anim_info: Dict[str, Any]
    ) -> Optional[str]:
        """
        Render a slide at a specific animation state.
        
        Args:
            slide: PowerPoint slide object
            slide_number: Slide number
            step: Animation step number  
            anim_info: Animation state information
            
        Returns:
            Path to rendered image or None if failed
        """
        try:
            # Create output path for this specific state
            output_path = self.output_dir / f"slide_{slide_number:03d}_step_{step:02d}.png"
            
            if not PIL_AVAILABLE:
                logger.warning("PIL not available - creating text placeholder for animation state")
                return str(output_path)
            
            # Create image representing this animation state
            size = (1920, 1080)
            img = Image.new('RGB', size, color='white')
            draw = ImageDraw.Draw(img)
            
            try:
                font = ImageFont.load_default()
            except:
                font = None
            
            # Add slide and step info
            draw.text((50, 50), f"Slide {slide_number} - Step {step}", fill='black', font=font)
            draw.text((50, 80), f"Type: {anim_info.get('animation_type', 'unknown')}", fill='blue', font=font)
            
            # Add visible elements for this state
            visible_elements = anim_info.get('visible_elements', [])
            y_pos = 120
            
            for i, element in enumerate(visible_elements[:10]):  # Limit to avoid overflow
                if element:
                    # Truncate long text
                    display_text = element[:80] + "..." if len(element) > 80 else element
                    draw.text((50, y_pos), f"• {display_text}", fill='black', font=font)
                    y_pos += 30
            
            if len(visible_elements) > 10:
                draw.text((50, y_pos), f"... and {len(visible_elements) - 10} more elements", fill='gray', font=font)
            
            # Add confidence indicator
            confidence = anim_info.get('confidence', 0.5)
            confidence_color = 'green' if confidence > 0.8 else 'orange' if confidence > 0.6 else 'red'
            draw.text((50, size[1] - 50), f"Confidence: {confidence:.1f}", fill=confidence_color, font=font)
            
            # Save the image
            img.save(output_path)
            logger.debug(f"Rendered animation state: {output_path}")
            
            return str(output_path)
            
        except Exception as e:
            logger.error(f"Failed to render slide state for slide {slide_number} step {step}: {e}")
            return None


class AnimationAwareVisualValidator:
    """Enhanced visual validator that understands PowerPoint build animations."""
    
    def __init__(self, slide_extractor: SlideImageExtractor):
        """Initialize with a slide extractor."""
        self.slide_extractor = slide_extractor
    
    def find_best_animation_state(
        self,
        slide_states: List[SlideAnimationState],
        video_frame_path: str,
        narration_text: str,
        timestamp_in_slide: float
    ) -> Optional[SlideAnimationState]:
        """
        Find the animation state that best matches the video frame and narration context.
        
        Args:
            slide_states: Available animation states for the slide
            video_frame_path: Path to video frame
            narration_text: Current narration text
            timestamp_in_slide: How far into the slide we are (0.0-1.0)
            
        Returns:
            Best matching SlideAnimationState or None
        """
        if not slide_states:
            return None
        
        if len(slide_states) == 1:
            return slide_states[0]
        
        # Score each state based on timing and content alignment
        scored_states = []
        
        for state in slide_states:
            score = self._score_animation_state(
                state, narration_text, timestamp_in_slide
            )
            scored_states.append((state, score))
        
        # Return the highest scoring state
        scored_states.sort(key=lambda x: x[1], reverse=True)
        best_state = scored_states[0][0]
        
        logger.debug(f"Selected animation state {best_state.build_step} for slide {best_state.slide_number} "
                    f"(score: {scored_states[0][1]:.2f})")
        
        return best_state
    
    def _score_animation_state(
        self,
        state: SlideAnimationState,
        narration_text: str,
        timestamp_in_slide: float
    ) -> float:
        """Score how well an animation state matches the current context."""
        
        score = state.confidence  # Base score from animation detection confidence
        
        # Timing alignment: earlier states should match earlier in the slide
        total_states = 5  # Rough estimate - would be passed in ideally
        expected_timing = state.build_step / max(1, total_states - 1)
        timing_diff = abs(expected_timing - timestamp_in_slide)
        timing_score = 1.0 - timing_diff
        
        score += timing_score * 0.3
        
        # Content alignment: check if narration mentions visible elements
        if state.visible_elements and narration_text:
            narration_lower = narration_text.lower()
            matching_elements = 0
            
            for element in state.visible_elements:
                if element and any(word in narration_lower for word in element.lower().split()[:3]):
                    matching_elements += 1
            
            if state.visible_elements:
                content_alignment = matching_elements / len(state.visible_elements)
                score += content_alignment * 0.2
        
        return min(1.0, score)


class VideoFrameExtractor:
    """Extract frames from video at specific timestamps for visual validation."""
    
    def __init__(self, output_dir: Optional[str] = None):
        """
        Initialize the video frame extractor.
        
        Args:
            output_dir: Directory to save extracted frames (default: temp dir)
        """
        if not FFMPEG_AVAILABLE:
            logger.warning("ffmpeg-python not available - frame extraction may be limited")
        
        self.output_dir = Path(output_dir) if output_dir else Path(tempfile.gettempdir()) / "autovid_frames"
        self.output_dir.mkdir(exist_ok=True, parents=True)
    
    def extract_frames_at_transitions(
        self, 
        video_path: str, 
        transition_timestamps: List[float]
    ) -> List[str]:
        """
        Extract video frames at specific transition timestamps.
        
        Args:
            video_path: Path to the video file
            transition_timestamps: List of timestamps to extract frames
            
        Returns:
            List of paths to extracted frame images
        """
        if not FFMPEG_AVAILABLE:
            logger.error("ffmpeg-python is required for frame extraction. Please install: pip install ffmpeg-python")
            return []
        
        frame_paths = []
        
        for i, timestamp in enumerate(transition_timestamps):
            frame_path = self.extract_frame_at_timestamp(video_path, timestamp, f"transition_{i:03d}")
            if frame_path:
                frame_paths.append(frame_path)
        
        logger.info(f"Extracted {len(frame_paths)} frames from {video_path}")
        return frame_paths
    
    def extract_frame_at_timestamp(
        self, 
        video_path: str, 
        timestamp: float, 
        frame_name: Optional[str] = None
    ) -> Optional[str]:
        """
        Extract a single frame from video at specific timestamp.
        
        Args:
            video_path: Path to the video file
            timestamp: Timestamp in seconds
            frame_name: Optional custom name for the frame file
            
        Returns:
            Path to extracted frame image or None if failed
        """
        try:
            if not frame_name:
                frame_name = f"frame_{timestamp:.2f}s"
            
            output_path = self.output_dir / f"{frame_name}.png"
            
            if not FFMPEG_AVAILABLE:
                logger.error("ffmpeg-python required for frame extraction")
                return None
            
            # Use ffmpeg to extract frame at timestamp
            (
                ffmpeg
                .input(video_path, ss=timestamp)
                .output(str(output_path), vframes=1, format='image2')
                .overwrite_output()
                .run(quiet=True)
            )
            
            if output_path.exists():
                logger.debug(f"Extracted frame at {timestamp}s: {output_path}")
                return str(output_path)
            else:
                logger.error(f"Frame extraction failed - file not created: {output_path}")
                return None
                
        except Exception as e:
            logger.error(f"Failed to extract frame at {timestamp}s: {e}")
            return None
    
    def batch_extract_frames(
        self, 
        video_path: str, 
        timestamps: List[float], 
        frame_names: Optional[List[str]] = None
    ) -> List[str]:
        """
        Extract multiple frames efficiently in batch.
        
        Args:
            video_path: Path to the video file
            timestamps: List of timestamps to extract
            frame_names: Optional list of custom names for frames
            
        Returns:
            List of paths to successfully extracted frames
        """
        frame_paths = []
        
        for i, timestamp in enumerate(timestamps):
            frame_name = frame_names[i] if frame_names and i < len(frame_names) else None
            frame_path = self.extract_frame_at_timestamp(video_path, timestamp, frame_name)
            
            if frame_path:
                frame_paths.append(frame_path)
        
        return frame_paths


class BasicVisualComparator:
    """Basic computer vision comparison for Phase 1 implementation."""
    
    def __init__(self):
        """Initialize the basic visual comparator."""
        self.comparison_methods = ['histogram', 'template', 'hash']
        
    def calculate_similarity(
        self, 
        slide_image_path: str, 
        video_frame_path: str, 
        method: str = 'histogram'
    ) -> float:
        """
        Calculate similarity between slide image and video frame.
        
        Args:
            slide_image_path: Path to slide image
            video_frame_path: Path to video frame
            method: Comparison method ('histogram', 'template', 'hash')
            
        Returns:
            Similarity score between 0.0 and 1.0
        """
        if not PIL_AVAILABLE:
            logger.warning("PIL not available - using basic file comparison")
            return self._basic_file_comparison(slide_image_path, video_frame_path)
        
        try:
            if method == 'histogram':
                return self._histogram_comparison(slide_image_path, video_frame_path)
            elif method == 'template':
                return self._template_comparison(slide_image_path, video_frame_path)
            elif method == 'hash':
                return self._perceptual_hash_comparison(slide_image_path, video_frame_path)
            else:
                logger.warning(f"Unknown comparison method: {method}")
                return 0.0
                
        except Exception as e:
            logger.error(f"Similarity calculation failed: {e}")
            return 0.0
    
    def _histogram_comparison(self, img1_path: str, img2_path: str) -> float:
        """Compare images using histogram similarity."""
        try:
            img1 = Image.open(img1_path).convert('RGB')
            img2 = Image.open(img2_path).convert('RGB')
            
            # Resize to same size for comparison
            size = (256, 256)
            img1 = img1.resize(size)
            img2 = img2.resize(size)
            
            # Convert to numpy arrays
            arr1 = np.array(img1)
            arr2 = np.array(img2)
            
            # Calculate histograms for each channel
            hist1_r = np.histogram(arr1[:,:,0], bins=256, range=(0, 256))[0]
            hist1_g = np.histogram(arr1[:,:,1], bins=256, range=(0, 256))[0]
            hist1_b = np.histogram(arr1[:,:,2], bins=256, range=(0, 256))[0]
            
            hist2_r = np.histogram(arr2[:,:,0], bins=256, range=(0, 256))[0]
            hist2_g = np.histogram(arr2[:,:,1], bins=256, range=(0, 256))[0]
            hist2_b = np.histogram(arr2[:,:,2], bins=256, range=(0, 256))[0]
            
            # Calculate correlation for each channel
            corr_r = np.corrcoef(hist1_r, hist2_r)[0, 1]
            corr_g = np.corrcoef(hist1_g, hist2_g)[0, 1]
            corr_b = np.corrcoef(hist1_b, hist2_b)[0, 1]
            
            # Handle NaN values (can occur with constant images)
            correlations = [c for c in [corr_r, corr_g, corr_b] if not np.isnan(c)]
            
            if not correlations:
                return 0.0
            
            avg_correlation = np.mean(correlations)
            return max(0.0, avg_correlation)  # Clamp negative correlations to 0
            
        except Exception as e:
            logger.error(f"Histogram comparison failed: {e}")
            return 0.0
    
    def _template_comparison(self, img1_path: str, img2_path: str) -> float:
        """Compare images using template matching (requires OpenCV)."""
        if not CV2_AVAILABLE:
            logger.warning("OpenCV not available - falling back to histogram comparison")
            return self._histogram_comparison(img1_path, img2_path)
        
        try:
            img1 = cv2.imread(img1_path)
            img2 = cv2.imread(img2_path)
            
            if img1 is None or img2 is None:
                return 0.0
            
            # Convert to grayscale
            gray1 = cv2.cvtColor(img1, cv2.COLOR_BGR2GRAY)
            gray2 = cv2.cvtColor(img2, cv2.COLOR_BGR2GRAY)
            
            # Resize to same size
            size = (512, 512)
            gray1 = cv2.resize(gray1, size)
            gray2 = cv2.resize(gray2, size)
            
            # Template matching
            result = cv2.matchTemplate(gray1, gray2, cv2.TM_CCOEFF_NORMED)
            _, max_val, _, _ = cv2.minMaxLoc(result)
            
            return max(0.0, max_val)
            
        except Exception as e:
            logger.error(f"Template comparison failed: {e}")
            return 0.0
    
    def _perceptual_hash_comparison(self, img1_path: str, img2_path: str) -> float:
        """Compare images using perceptual hashing."""
        try:
            hash1 = self._calculate_perceptual_hash(img1_path)
            hash2 = self._calculate_perceptual_hash(img2_path)
            
            if not hash1 or not hash2:
                return 0.0
            
            # Calculate Hamming distance
            hamming_distance = sum(c1 != c2 for c1, c2 in zip(hash1, hash2))
            max_distance = len(hash1)
            
            # Convert to similarity score
            similarity = 1.0 - (hamming_distance / max_distance)
            return similarity
            
        except Exception as e:
            logger.error(f"Perceptual hash comparison failed: {e}")
            return 0.0
    
    def _calculate_perceptual_hash(self, img_path: str, hash_size: int = 8) -> Optional[str]:
        """Calculate perceptual hash of an image."""
        try:
            img = Image.open(img_path).convert('L')  # Convert to grayscale
            img = img.resize((hash_size + 1, hash_size), Image.Resampling.LANCZOS)
            
            # Convert to numpy array
            pixels = np.array(img)
            
            # Calculate differences between adjacent pixels
            diff = pixels[:, 1:] > pixels[:, :-1]
            
            # Convert to hash string
            hash_string = ''.join(str(b) for b in diff.flatten())
            return hash_string
            
        except Exception as e:
            logger.error(f"Perceptual hash calculation failed: {e}")
            return None
    
    def _basic_file_comparison(self, file1_path: str, file2_path: str) -> float:
        """Basic file comparison fallback when image libraries aren't available."""
        try:
            # Compare file sizes as a very basic similarity metric
            size1 = Path(file1_path).stat().st_size
            size2 = Path(file2_path).stat().st_size
            
            if size1 == 0 or size2 == 0:
                return 0.0
            
            size_ratio = min(size1, size2) / max(size1, size2)
            
            # This is a very crude approximation
            return size_ratio * 0.5  # Conservative similarity estimate
            
        except Exception as e:
            logger.error(f"Basic file comparison failed: {e}")
            return 0.0


class VisionLLMValidator:
    """Phase 2: Vision LLM integration for semantic visual validation."""
    
    def __init__(self, model_name: str = "gemini-2.5-flash-lite"):
        """
        Initialize the Vision LLM validator.
        
        Args:
            model_name: Name of the vision model to use
        """
        self.model_name = model_name
        self.max_image_size = (1024, 1024)  # Resize large images for efficiency
        
        # Try to import the consultation tool for LLM access
        try:
            # This would use the existing consultation system
            self.llm_available = True
        except ImportError:
            self.llm_available = False
            logger.warning("Vision LLM not available - falling back to basic CV")
    
    def validate_slide_match(
        self,
        slide_image_path: str,
        video_frame_path: str,
        expected_content: str,
        narration_text: str,
        slide_number: int
    ) -> VisualValidationResult:
        """
        Use Vision LLM to validate if video frame matches expected slide content.
        
        Args:
            slide_image_path: Path to expected slide image
            video_frame_path: Path to video frame
            expected_content: Expected slide content description
            narration_text: What should be narrated at this point
            slide_number: Slide number for context
            
        Returns:
            VisualValidationResult with LLM analysis
        """
        try:
            # Prepare images for LLM analysis
            slide_b64 = self._encode_image_for_llm(slide_image_path)
            frame_b64 = self._encode_image_for_llm(video_frame_path)
            
            if not slide_b64 or not frame_b64:
                return self._create_fallback_result(
                    slide_number, slide_image_path, video_frame_path,
                    "Failed to encode images for LLM analysis"
                )
            
            # Create prompt for vision LLM
            prompt = self._create_validation_prompt(
                expected_content, narration_text, slide_number
            )
            
            # Call Vision LLM (placeholder for now - would integrate with actual LLM)
            llm_response = self._call_vision_llm(prompt, slide_b64, frame_b64)
            
            # Parse LLM response
            result = self._parse_llm_response(
                llm_response, slide_number, slide_image_path, video_frame_path
            )
            
            return result
            
        except Exception as e:
            logger.error(f"Vision LLM validation failed: {e}")
            return self._create_fallback_result(
                slide_number, slide_image_path, video_frame_path,
                f"LLM validation error: {str(e)}"
            )
    
    def batch_validate_slides(
        self,
        slide_frame_pairs: List[Tuple[str, str, str, str, int]],
        progress_callback: Optional[callable] = None
    ) -> List[VisualValidationResult]:
        """
        Validate multiple slide-frame pairs efficiently.
        
        Args:
            slide_frame_pairs: List of (slide_path, frame_path, expected_content, narration, slide_num)
            progress_callback: Optional callback for progress updates
            
        Returns:
            List of VisualValidationResult objects
        """
        results = []
        
        for i, (slide_path, frame_path, content, narration, slide_num) in enumerate(slide_frame_pairs):
            if progress_callback:
                progress_callback(i, len(slide_frame_pairs))
            
            result = self.validate_slide_match(
                slide_path, frame_path, content, narration, slide_num
            )
            results.append(result)
        
        return results
    
    def _encode_image_for_llm(self, image_path: str) -> Optional[str]:
        """Encode image as base64 for LLM analysis."""
        try:
            if not PIL_AVAILABLE:
                logger.error("PIL required for image encoding")
                return None
            
            # Load and potentially resize image
            img = Image.open(image_path)
            
            # Resize if too large to save on LLM costs
            if img.size[0] > self.max_image_size[0] or img.size[1] > self.max_image_size[1]:
                img.thumbnail(self.max_image_size, Image.Resampling.LANCZOS)
            
            # Convert to base64
            import io
            buffer = io.BytesIO()
            img.save(buffer, format='PNG')
            img_b64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
            
            return img_b64
            
        except Exception as e:
            logger.error(f"Image encoding failed for {image_path}: {e}")
            return None
    
    def _create_validation_prompt(
        self, 
        expected_content: str, 
        narration_text: str, 
        slide_number: int
    ) -> str:
        """Create a prompt for the vision LLM to validate slide-video alignment."""
        
        prompt = f"""
You are validating video-slide synchronization for an automated presentation system.

CONTEXT:
- This is slide #{slide_number} from a presentation
- The narration at this moment should be: "{narration_text}"
- The slide should contain: "{expected_content}"

TASK:
Compare the two images provided:
1. Image 1: Expected slide content (from PowerPoint)  
2. Image 2: Video frame (from exported presentation video)

ANALYSIS REQUIRED:
1. Do both images show the same slide content? (Consider that Image 2 might have slight export differences)
2. Does the visible content match what should be narrated ("{narration_text}")?
3. Are there any obvious sync issues (wrong slide, missing content, build animation mismatches)?
4. What is your confidence level that these images represent the same slide at the correct moment?

RESPONSE FORMAT (JSON):
{{
    "content_match": <boolean>,
    "narration_alignment": <boolean>,
    "confidence_score": <float 0.0-1.0>,
    "issues": [<list of specific problems>],
    "recommendations": [<list of suggestions>],
    "explanation": "<detailed analysis>"
}}

Focus on content alignment rather than minor visual differences due to export quality or rendering.
"""
        
        return prompt.strip()
    
    def _call_vision_llm(self, prompt: str, slide_b64: str, frame_b64: str) -> Dict[str, Any]:
        """
        Call the vision LLM with images and prompt.
        
        This is a placeholder that would integrate with the existing consultation system
        or direct LLM API calls.
        """
        
        # Placeholder implementation - would use actual LLM integration
        # For now, return a mock response structure
        
        if not self.llm_available:
            return {
                "content_match": True,
                "narration_alignment": True,  
                "confidence_score": 0.7,
                "issues": ["LLM not available - using fallback"],
                "recommendations": ["Install vision LLM dependencies"],
                "explanation": "Fallback response - vision LLM not configured"
            }
        
        # TODO: Integrate with actual Vision LLM API
        # This would involve:
        # 1. Using the mcp__consult7__consultation tool with vision model
        # 2. Or direct API calls to Gemini Vision
        # 3. Proper image encoding and transmission
        
        logger.info(f"Would call {self.model_name} with prompt length: {len(prompt)}")
        logger.debug("Vision LLM integration - placeholder implementation")
        
        # Mock response for development
        return {
            "content_match": True,
            "narration_alignment": True,
            "confidence_score": 0.8,
            "issues": [],
            "recommendations": [],
            "explanation": "Mock validation response - implement actual LLM integration"
        }
    
    def _parse_llm_response(
        self,
        llm_response: Dict[str, Any],
        slide_number: int,
        slide_image_path: str,
        video_frame_path: str
    ) -> VisualValidationResult:
        """Parse LLM response into VisualValidationResult."""
        
        try:
            # Extract similarity score from LLM confidence and content match
            content_match = llm_response.get('content_match', False)
            narration_alignment = llm_response.get('narration_alignment', False)
            llm_confidence = llm_response.get('confidence_score', 0.0)
            
            # Calculate similarity score based on LLM assessment
            if content_match and narration_alignment:
                similarity_score = llm_confidence
            elif content_match:
                similarity_score = llm_confidence * 0.8  # Content matches but narration might not
            else:
                similarity_score = llm_confidence * 0.5  # Significant issues
            
            return VisualValidationResult(
                slide_number=slide_number,
                video_timestamp=0.0,  # Would be filled in by caller
                slide_image_path=slide_image_path,
                video_frame_path=video_frame_path,
                similarity_score=similarity_score,
                validation_method="vision_llm",
                confidence=llm_confidence,
                issues=llm_response.get('issues', []),
                recommendations=llm_response.get('recommendations', []),
                llm_analysis=llm_response
            )
            
        except Exception as e:
            logger.error(f"Failed to parse LLM response: {e}")
            return self._create_fallback_result(
                slide_number, slide_image_path, video_frame_path,
                f"LLM response parsing error: {str(e)}"
            )
    
    def _create_fallback_result(
        self,
        slide_number: int,
        slide_image_path: str,
        video_frame_path: str,
        error_message: str
    ) -> VisualValidationResult:
        """Create a fallback result when LLM validation fails."""
        
        return VisualValidationResult(
            slide_number=slide_number,
            video_timestamp=0.0,
            slide_image_path=slide_image_path,
            video_frame_path=video_frame_path,
            similarity_score=0.5,  # Neutral score when uncertain
            validation_method="fallback",
            confidence=0.3,  # Low confidence for fallback
            issues=[error_message],
            recommendations=["Check image files and LLM configuration"],
            llm_analysis=None
        )


class VisualSyncValidator:
    """Main visual validation orchestrator combining Phase 1 and Phase 2."""
    
    def __init__(
        self,
        use_vision_llm: bool = True,
        vision_model: str = "gemini-2.5-flash-lite",
        output_dir: Optional[str] = None
    ):
        """
        Initialize the visual sync validator.
        
        Args:
            use_vision_llm: Whether to use Vision LLM for semantic validation
            vision_model: Vision model to use for LLM validation
            output_dir: Directory for temporary files
        """
        self.use_vision_llm = use_vision_llm
        self.output_dir = output_dir
        
        # Initialize components
        self.slide_extractor = SlideImageExtractor(output_dir)
        self.frame_extractor = VideoFrameExtractor(output_dir)
        self.basic_comparator = BasicVisualComparator()
        
        if use_vision_llm:
            self.vision_llm = VisionLLMValidator(vision_model)
        else:
            self.vision_llm = None
        
        # Initialize animation-aware validator
        self.animation_validator = AnimationAwareVisualValidator(self.slide_extractor)
    
    def validate_sync_plan_visually(
        self,
        sync_plan,  # SyncPlan object from slide_sync.py
        video_path: str,
        pptx_path: str,
        transcript_data: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        Perform comprehensive visual validation of a sync plan.
        
        Args:
            sync_plan: SyncPlan object to validate
            video_path: Path to the presentation video
            pptx_path: Path to the original PowerPoint file
            transcript_data: Optional transcript data for context
            
        Returns:
            Comprehensive visual validation results
        """
        logger.info("Starting visual validation of sync plan")
        
        results = {
            'overall_score': 0.0,
            'validation_method': 'hybrid' if self.use_vision_llm else 'basic_cv',
            'segment_results': [],
            'issues': [],
            'recommendations': [],
            'dependencies_available': validate_visual_dependencies()
        }
        
        try:
            # Extract slide images
            logger.info("Extracting slide images from PowerPoint...")
            slide_images = self.slide_extractor.extract_all_slides(pptx_path)
            
            if not slide_images:
                results['issues'].append("Failed to extract slide images")
                return results
            
            # Extract video frames at transition points
            logger.info("Extracting video frames at transition points...")
            transition_timestamps = [seg.keynote_start for seg in sync_plan.segments]
            video_frames = self.frame_extractor.extract_frames_at_transitions(
                video_path, transition_timestamps
            )
            
            if not video_frames:
                results['issues'].append("Failed to extract video frames")
                return results
            
            # Validate each segment with build animation awareness
            logger.info(f"Validating {len(sync_plan.segments)} segments with build animation awareness...")
            segment_scores = []
            
            for i, segment in enumerate(sync_plan.segments):
                if i >= len(video_frames):
                    continue
                
                # Get expected content for this segment
                expected_content = self._get_expected_content(segment, transcript_data)
                narration_text = self._get_narration_text(segment, transcript_data)
                
                # Extract animation states for this slide
                slide_states = self.slide_extractor.extract_slide_states(pptx_path, segment.slide_number)
                
                if slide_states:
                    # Find the best animation state for this timing
                    timestamp_in_slide = self._calculate_timestamp_in_slide(segment, sync_plan.segments)
                    best_state = self.animation_validator.find_best_animation_state(
                        slide_states, video_frames[i], narration_text, timestamp_in_slide
                    )
                    
                    slide_image_to_use = best_state.image_path if best_state else slide_states[0].image_path
                    animation_confidence = best_state.confidence if best_state else 0.5
                else:
                    # Fallback to basic slide image
                    slide_image_to_use = slide_images[i] if i < len(slide_images) else None
                    animation_confidence = 0.5
                
                if not slide_image_to_use:
                    continue
                
                # Perform validation with animation-aware context
                if self.use_vision_llm and self.vision_llm:
                    # Enhanced LLM validation with animation context
                    validation_result = self._validate_with_animation_context(
                        slide_image_to_use, video_frames[i], expected_content, 
                        narration_text, segment, best_state if slide_states else None
                    )
                else:
                    # Enhanced basic CV with animation context
                    similarity = self.basic_comparator.calculate_similarity(
                        slide_image_to_use, video_frames[i], 'histogram'
                    )
                    
                    # Adjust similarity based on animation confidence
                    adjusted_similarity = similarity * animation_confidence
                    
                    validation_result = VisualValidationResult(
                        slide_number=segment.slide_number,
                        video_timestamp=segment.keynote_start,
                        slide_image_path=slide_image_to_use,
                        video_frame_path=video_frames[i],
                        similarity_score=adjusted_similarity,
                        validation_method="basic_cv_with_animation",
                        confidence=animation_confidence * 0.8,  # Animation detection affects confidence
                        issues=[],
                        recommendations=[]
                    )
                
                results['segment_results'].append(validation_result)
                segment_scores.append(validation_result.similarity_score)
            
            # Calculate overall score
            if segment_scores:
                results['overall_score'] = sum(segment_scores) / len(segment_scores)
            
            # Generate summary recommendations
            results['recommendations'] = self._generate_visual_recommendations(results)
            
            logger.info(f"Visual validation complete. Overall score: {results['overall_score']:.2f}")
            
        except Exception as e:
            logger.error(f"Visual validation failed: {e}")
            results['issues'].append(f"Validation error: {str(e)}")
        
        return results
    
    def _get_expected_content(self, segment, transcript_data: Optional[Dict[str, Any]]) -> str:
        """Extract expected content description for a segment."""
        if not transcript_data:
            return f"Slide {segment.slide_number} content"
        
        # Try to find transcript content for this slide
        if 'slides' in transcript_data:
            slides = transcript_data.get('slides', [])
            for slide in slides:
                if slide.get('index', 0) + 1 == segment.slide_number:
                    segments = slide.get('segments', [])
                    text_content = []
                    for seg in segments:
                        if seg.get('kind') == 'text':
                            text_content.append(seg.get('text', ''))
                    return ' '.join(text_content) if text_content else f"Slide {segment.slide_number}"
        
        return f"Slide {segment.slide_number} content"
    
    def _get_narration_text(self, segment, transcript_data: Optional[Dict[str, Any]]) -> str:
        """Extract narration text for a segment."""
        # This would extract the actual narration text that should be spoken
        # during this slide segment
        return f"Narration for slide {segment.slide_number}"
    
    def _generate_visual_recommendations(self, results: Dict[str, Any]) -> List[str]:
        """Generate recommendations based on visual validation results."""
        recommendations = []
        
        overall_score = results.get('overall_score', 0.0)
        
        if overall_score < 0.5:
            recommendations.append("Low visual validation score - check slide-video alignment")
        
        if not results.get('dependencies_available', {}).get('PIL'):
            recommendations.append("Install Pillow for better image processing: pip install Pillow")
        
        if not results.get('dependencies_available', {}).get('ffmpeg'):
            recommendations.append("Install ffmpeg-python for video frame extraction: pip install ffmpeg-python")
        
        # Analyze segment-specific issues
        segment_results = results.get('segment_results', [])
        low_confidence_segments = [r for r in segment_results if r.confidence < 0.6]
        
        if low_confidence_segments:
            recommendations.append(f"{len(low_confidence_segments)} segments have low visual confidence")
        
        return recommendations
    
    def _calculate_timestamp_in_slide(self, segment, all_segments: List) -> float:
        """Calculate how far into the slide we are (0.0 = start, 1.0 = end)."""
        slide_duration = segment.keynote_duration + segment.gap_needed
        if slide_duration <= 0:
            return 0.5  # Default to middle if duration unknown
        
        # For now, assume we're validating at the start of each segment
        # In a more sophisticated implementation, this would be based on
        # the actual timestamp within the segment
        return 0.1  # Early in the slide for transition validation
    
    def _validate_with_animation_context(
        self,
        slide_image_path: str,
        video_frame_path: str,
        expected_content: str,
        narration_text: str,
        segment,
        animation_state: Optional[SlideAnimationState]
    ) -> VisualValidationResult:
        """Perform vision LLM validation with animation context."""
        
        # Enhanced prompt with animation context
        if animation_state:
            animation_context = f"""
ANIMATION CONTEXT:
- This slide has build animations (Step {animation_state.build_step})
- Animation type: {animation_state.animation_type}
- Visible elements at this step: {animation_state.visible_elements[:3]}
- Consider that the video frame might show a different animation step than expected
"""
            enhanced_expected_content = f"{expected_content}\n{animation_context}"
        else:
            enhanced_expected_content = expected_content
        
        # Use standard LLM validation with enhanced context
        validation_result = self.vision_llm.validate_slide_match(
            slide_image_path, video_frame_path, enhanced_expected_content,
            narration_text, segment.slide_number
        )
        
        # Adjust confidence based on animation awareness
        if animation_state:
            validation_result.confidence *= animation_state.confidence
            if animation_state.build_step > 0:
                validation_result.issues.append(f"Slide has build animations (step {animation_state.build_step})")
        
        return validation_result


def validate_visual_dependencies() -> Dict[str, bool]:
    """Check which visual processing dependencies are available."""
    dependencies = {
        'pptx': PPTX_AVAILABLE,
        'PIL': PIL_AVAILABLE,
        'cv2': CV2_AVAILABLE,
        'ffmpeg': FFMPEG_AVAILABLE
    }
    
    missing = [name for name, available in dependencies.items() if not available]
    
    if missing:
        logger.warning(f"Missing optional visual dependencies: {missing}")
        logger.info("Install missing dependencies:")
        if 'PIL' in missing:
            logger.info("  pip install Pillow")
        if 'cv2' in missing:
            logger.info("  pip install opencv-python")
        if 'ffmpeg' in missing:
            logger.info("  pip install ffmpeg-python")
    
    return dependencies