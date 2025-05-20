import pytest
import subprocess
import json
import os
from unittest.mock import patch # Add this import
from pptx import Presentation
from cli.extract_transcript import extract_transcript # Import the function directly for some tests

# Define the path for the dummy PPTX file
DUMMY_PPTX_PATH = "tests/fixtures/dummy_presentation_for_cues.pptx"
OUTPUT_MANIFEST_PATH_CLI = "workspace/01_notes/dummy_cue_manifest_cli.json"

# Expected default cue token
DEFAULT_CUE = "[transition]"

@pytest.fixture(scope="module")
def create_dummy_pptx_for_cues():
    """Creates a dummy PPTX file with various speaker notes for cue testing."""
    prs = Presentation()
    # Ensure the fixtures directory exists
    os.makedirs(os.path.dirname(DUMMY_PPTX_PATH), exist_ok=True)

    # Slide 0: No cues, just text
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 0: No Cues"
    slide.notes_slide.notes_text_frame.text = "This is simple text without any cues."

    # Slide 1: Default cue token "[transition]" present
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 1: Default Cue"
    slide.notes_slide.notes_text_frame.text = f"Part 1 {DEFAULT_CUE} Part 2"

    # Slide 2: Default cue token, case variation "[Transition]"
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 2: Default Cue - Case Variation 1"
    slide.notes_slide.notes_text_frame.text = "Text A [Transition] Text B"

    # Slide 3: Default cue token, all caps "[TRANSITION]"
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 3: Default Cue - Case Variation 2"
    slide.notes_slide.notes_text_frame.text = "Segment X [TRANSITION] Segment Y"

    # Slide 4: Multiple default cues
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 4: Multiple Default Cues"
    slide.notes_slide.notes_text_frame.text = f"Start {DEFAULT_CUE} Middle {DEFAULT_CUE} End"
    
    # Slide 5: Cue at the beginning
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 5: Cue at Beginning"
    slide.notes_slide.notes_text_frame.text = f"{DEFAULT_CUE} Text after cue"

    # Slide 6: Cue at the end
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 6: Cue at End"
    slide.notes_slide.notes_text_frame.text = f"Text before cue {DEFAULT_CUE}"

    # Slide 7: Cue with special regex characters (e.g., "*")
    # For this, we'll test with a custom cue token in the test itself.
    # Notes: "Alpha *CUE* Beta"
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 7: Special Char Cue"
    slide.notes_slide.notes_text_frame.text = "Alpha *CUE* Beta"

    # Slide 8: Text that looks like a cue but isn't the defined one
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 8: False Positive Cue"
    slide.notes_slide.notes_text_frame.text = f"This is text with transition (no brackets) and also {DEFAULT_CUE} a real cue."

    # Slide 9: Ellipses interaction
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 9: Ellipses and Cue"
    slide.notes_slide.notes_text_frame.text = f"...Start text... {DEFAULT_CUE} …Middle text… {DEFAULT_CUE} ...End text..."

    # Slide 10: Empty notes
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 10: Empty Notes"
    slide.notes_slide.notes_text_frame.text = ""
    
    # Slide 11: Notes with only whitespace
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 11: Whitespace Notes"
    slide.notes_slide.notes_text_frame.text = "   \n\t   "

    # Slide 12: Unicode ellipses and cues
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Slide 12: Unicode Ellipses and Cues"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"First part… {DEFAULT_CUE} …Second part"
    
    # Slide 13: Only a cue token
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Slide 13: Only a cue"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"{DEFAULT_CUE}"

    # Slide 14: Cue with leading/trailing whitespace in notes
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Slide 14: Cue with surrounding whitespace"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"Text before \t {DEFAULT_CUE} \n Text after"


    prs.save(DUMMY_PPTX_PATH)
    yield DUMMY_PPTX_PATH

    if os.path.exists(DUMMY_PPTX_PATH):
        os.remove(DUMMY_PPTX_PATH)
    if os.path.exists(OUTPUT_MANIFEST_PATH_CLI):
        os.remove(OUTPUT_MANIFEST_PATH_CLI)
    output_dir_cli = os.path.dirname(OUTPUT_MANIFEST_PATH_CLI)
    if os.path.exists(output_dir_cli) and not os.listdir(output_dir_cli):
        os.rmdir(output_dir_cli)

def run_script(pptx_path, out_path, cue_token, step_id="test_step"):
    script_path = "cli/extract_transcript.py" # Corrected path
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    
    command = [
        "python", script_path,
        "--pptx", pptx_path,
        "--out", out_path,
        "--cue", cue_token,
        "--step_id", step_id
    ]
    result = subprocess.run(command, capture_output=True, text=True, check=False)
    
    # Try to load JSON from stdout if return code is non-zero, or from file if zero
    # This helps debug if the script fails before writing the file.
    manifest = None
    if result.stdout:
        try:
            manifest = json.loads(result.stdout)
        except json.JSONDecodeError:
            print("Failed to decode JSON from stdout:")
            print(result.stdout)
            print("Stderr:")
            print(result.stderr)


    if result.returncode == 0 and os.path.exists(out_path):
         with open(out_path, "r") as f:
            manifest = json.load(f)
    elif manifest is None : # if stdout was not valid JSON either
        print(f"Script execution failed or did not produce valid JSON. Return code: {result.returncode}")
        print("Stdout:")
        print(result.stdout)
        print("Stderr:")
        print(result.stderr)


    return result, manifest

# --- CLI Tests ---
def test_cli_default_cue_various_cases(create_dummy_pptx_for_cues):
    pptx_path = create_dummy_pptx_for_cues
    result, manifest = run_script(pptx_path, OUTPUT_MANIFEST_PATH_CLI, DEFAULT_CUE)

    assert result.returncode == 0, f"Script failed: {result.stderr}"
    assert manifest["status"] == "success"
    assert manifest["slide_count"] == 15

    # Slide 0: No cues
    assert manifest["slides"][0]["segments"] == [{"kind": "text", "text": "This is simple text without any cues."}]
    # Slide 1: Default cue "[transition]"
    assert manifest["slides"][1]["segments"] == [
        {"kind": "text", "text": "Part 1"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "Part 2"}
    ]
    # Slide 2: Default cue, case variation "[Transition]"
    assert manifest["slides"][2]["segments"] == [
        {"kind": "text", "text": "Text A"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "Text B"}
    ]
    # Slide 3: Default cue, all caps "[TRANSITION]"
    assert manifest["slides"][3]["segments"] == [
        {"kind": "text", "text": "Segment X"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "Segment Y"}
    ]
    # Slide 4: Multiple default cues
    assert manifest["slides"][4]["segments"] == [
        {"kind": "text", "text": "Start"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "Middle"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "End"}
    ]
    # Slide 5: Cue at beginning
    assert manifest["slides"][5]["segments"] == [
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "Text after cue"}
    ]
    # Slide 6: Cue at end
    assert manifest["slides"][6]["segments"] == [
        {"kind": "text", "text": "Text before cue"},
        {"kind": "cue", "cue": DEFAULT_CUE}
    ]
    # Slide 8: False positive cue
    assert manifest["slides"][8]["segments"] == [
        {"kind": "text", "text": "This is text with transition (no brackets) and also"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "a real cue."}
    ]
    # Slide 9: Ellipses and Cue
    assert manifest["slides"][9]["segments"] == [
        {"kind": "text", "text": "Start text", "continue": ["start", "end"]},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "Middle text", "continue": ["start", "end"]},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "End text", "continue": ["start", "end"]}
    ]
    # Slide 10: Empty notes
    assert manifest["slides"][10]["segments"] == []
    # Slide 11: Whitespace notes
    assert manifest["slides"][11]["segments"] == []
    # Slide 12: Unicode ellipses and cues
    assert manifest["slides"][12]["segments"] == [
        {"kind": "text", "text": "First part", "continue": "end"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "Second part", "continue": "start"}
    ]
    # Slide 13: Only a cue
    assert manifest["slides"][13]["segments"] == [{"kind": "cue", "cue": DEFAULT_CUE}]
    
    # Slide 14: Cue with surrounding whitespace
    assert manifest["slides"][14]["segments"] == [
        {"kind": "text", "text": "Text before"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "Text after"}
    ]


def test_cli_custom_cue_special_chars(create_dummy_pptx_for_cues):
    pptx_path = create_dummy_pptx_for_cues
    custom_cue = "*CUE*"
    result, manifest = run_script(pptx_path, OUTPUT_MANIFEST_PATH_CLI, custom_cue)

    assert result.returncode == 0, f"Script failed: {result.stderr}"
    assert manifest["status"] == "success"
    
    # Slide 7: Custom cue "*CUE*"
    assert manifest["slides"][7]["segments"] == [
        {"kind": "text", "text": "Alpha"},
        {"kind": "cue", "cue": custom_cue},
        {"kind": "text", "text": "Beta"}
    ]

def test_cli_no_notes_slide_allowed(create_dummy_pptx_for_cues):
    """Tests that slides with no notes are handled gracefully (empty segments)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # A layout that typically has a notes placeholder
    slide.shapes.title.text = "Slide with genuinely no notes field"
    # Deliberately do not access slide.notes_slide to ensure it's not created
    
    # Create a new slide that *does* have a notes_slide but notes_text_frame is None or empty
    slide_with_empty_notes_text_frame = prs.slides.add_slide(prs.slide_layouts[5])
    slide_with_empty_notes_text_frame.shapes.title.text = "Slide with empty notes_text_frame"
    # Access notes_slide, but don't add text. Depending on pptx library, notes_text_frame might be None.
    _ = slide_with_empty_notes_text_frame.notes_slide 
    # To be certain notes_text_frame.text is empty:
    slide_with_empty_notes_text_frame.notes_slide.notes_text_frame.text = ""


    no_notes_pptx_path = "tests/fixtures/no_notes_for_graceful_handling.pptx"
    prs.save(no_notes_pptx_path)

    result, manifest = run_script(no_notes_pptx_path, OUTPUT_MANIFEST_PATH_CLI, DEFAULT_CUE)
    
    assert result.returncode == 0, f"Script failed: {result.stderr}"
    assert manifest["status"] == "success"
    assert len(manifest["slides"]) == 2
    assert manifest["slides"][0]["segments"] == [] # Slide with no notes_slide or empty notes
    assert manifest["slides"][1]["segments"] == [] # Slide with notes_slide but empty text

    if os.path.exists(no_notes_pptx_path):
        os.remove(no_notes_pptx_path)

def test_cli_non_existent_pptx():
    non_existent_path = "tests/fixtures/non_existent.pptx"
    result, manifest = run_script(non_existent_path, OUTPUT_MANIFEST_PATH_CLI, DEFAULT_CUE)
    
    assert result.returncode == 1
    assert manifest["status"] == "failure"
    # python-pptx raises PackageNotFoundError for non-existent or malformed files
    assert "Package not found" in manifest["error"] or "No such file or directory" in manifest["error"]


# --- Direct Function Call Tests (for more granular testing if needed) ---

@pytest.fixture
def mock_slide_with_notes():
    class MockNotesTextFrame:
        def __init__(self, text):
            self.text = text

    class MockNotesSlide:
        def __init__(self, text=""):
            if text is None: # Simulate notes_text_frame being None
                self.notes_text_frame = None
            else:
                self.notes_text_frame = MockNotesTextFrame(text)

    class MockSlide:
        def __init__(self, notes_text=""):
            self.notes_slide = MockNotesSlide(notes_text)
            self.slide_id = 256 # Dummy id
            self.name = "Mock Slide"


    class MockPresentation:
        def __init__(self, slides_notes):
            self.slides = [MockSlide(notes) for notes in slides_notes]

    return MockPresentation

@patch('cli.extract_transcript.Presentation')
def test_direct_empty_notes_handling(mock_pptx_presentation, mock_slide_with_notes):
    # Test with notes_text_frame.text being empty
    mock_prs_instance = mock_slide_with_notes([""])
    mock_pptx_presentation.return_value = mock_prs_instance # When Presentation(mock_prs_instance) is called, return it.

    # Pass the mock_prs_instance, extract_transcript will internally call Presentation() which is now patched
    manifest = extract_transcript(mock_prs_instance) 
    assert manifest["status"] == "success"
    assert len(manifest["slides"]) == 1
    assert manifest["slides"][0]["segments"] == []

    # Test with notes_text_frame being None (if pptx library could return this)
    # Our current mock setup doesn't easily do this without more complex mocking or direct instantiation
    # For now, the "" test covers the practical case of empty notes.
    # If notes_slide.notes_text_frame itself can be None:
    mock_prs_instance_none_frame = mock_slide_with_notes([None])
    mock_pptx_presentation.return_value = mock_prs_instance_none_frame
    manifest_none = extract_transcript(mock_prs_instance_none_frame)
    assert manifest_none["status"] == "success"
    assert len(manifest_none["slides"]) == 1
    assert manifest_none["slides"][0]["segments"] == []

@patch('cli.extract_transcript.Presentation')
def test_direct_various_cue_scenarios(mock_pptx_presentation, mock_slide_with_notes):
    # Scenario 1: Basic cue
    mock_prs_instance = mock_slide_with_notes([f"Hello {DEFAULT_CUE} World"])
    mock_pptx_presentation.return_value = mock_prs_instance
    manifest = extract_transcript(mock_prs_instance, cue_token=DEFAULT_CUE)
    assert manifest["slides"][0]["segments"] == [
        {"kind": "text", "text": "Hello"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "World"}
    ]

    # Scenario 2: Case-insensitive cue
    mock_prs_instance = mock_slide_with_notes(["Hello [TRANSITION] World"])
    mock_pptx_presentation.return_value = mock_prs_instance
    manifest = extract_transcript(mock_prs_instance, cue_token=DEFAULT_CUE)
    assert manifest["slides"][0]["segments"] == [
        {"kind": "text", "text": "Hello"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "World"}
    ]

    # Scenario 3: Cue at start
    mock_prs_instance = mock_slide_with_notes([f"{DEFAULT_CUE}World"])
    mock_pptx_presentation.return_value = mock_prs_instance
    manifest = extract_transcript(mock_prs_instance, cue_token=DEFAULT_CUE)
    assert manifest["slides"][0]["segments"] == [
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "World"}
    ]

    # Scenario 4: Cue at end
    mock_prs_instance = mock_slide_with_notes([f"Hello{DEFAULT_CUE}"])
    mock_pptx_presentation.return_value = mock_prs_instance
    manifest = extract_transcript(mock_prs_instance, cue_token=DEFAULT_CUE)
    assert manifest["slides"][0]["segments"] == [
        {"kind": "text", "text": "Hello"},
        {"kind": "cue", "cue": DEFAULT_CUE}
    ]
    
    # Scenario 5: Only cue
    mock_prs_instance = mock_slide_with_notes([f"{DEFAULT_CUE}"])
    mock_pptx_presentation.return_value = mock_prs_instance
    manifest = extract_transcript(mock_prs_instance, cue_token=DEFAULT_CUE)
    assert manifest["slides"][0]["segments"] == [{"kind": "cue", "cue": DEFAULT_CUE}]

    # Scenario 6: Multiple cues
    mock_prs_instance = mock_slide_with_notes([f"A {DEFAULT_CUE} B {DEFAULT_CUE} C"])
    mock_pptx_presentation.return_value = mock_prs_instance
    manifest = extract_transcript(mock_prs_instance, cue_token=DEFAULT_CUE)
    assert manifest["slides"][0]["segments"] == [
        {"kind": "text", "text": "A"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "B"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "C"}
    ]

    # Scenario 7: Special character in cue
    custom_cue = "*STAR*"
    mock_prs_instance = mock_slide_with_notes([f"A {custom_cue} B"])
    mock_pptx_presentation.return_value = mock_prs_instance
    manifest = extract_transcript(mock_prs_instance, cue_token=custom_cue)
    assert manifest["slides"][0]["segments"] == [
        {"kind": "text", "text": "A"},
        {"kind": "cue", "cue": custom_cue},
        {"kind": "text", "text": "B"}
    ]

    # Scenario 8: No cues
    mock_prs_instance = mock_slide_with_notes(["Just plain text."])
    mock_pptx_presentation.return_value = mock_prs_instance
    manifest = extract_transcript(mock_prs_instance, cue_token=DEFAULT_CUE)
    assert manifest["slides"][0]["segments"] == [{"kind": "text", "text": "Just plain text."}]

    # Scenario 9: Ellipses and cues
    mock_prs_instance = mock_slide_with_notes([f"...Text1... {DEFAULT_CUE} …Text2…"])
    mock_pptx_presentation.return_value = mock_prs_instance
    manifest = extract_transcript(mock_prs_instance, cue_token=DEFAULT_CUE)
    assert manifest["slides"][0]["segments"] == [
        {"kind": "text", "text": "Text1", "continue": ["start", "end"]},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "Text2", "continue": ["start", "end"]}
    ]

    # Scenario 10: Text resembling cue but not matching
    mock_prs_instance = mock_slide_with_notes([f"This is [some text] not a cue. But this is {DEFAULT_CUE} a cue."])
    mock_pptx_presentation.return_value = mock_prs_instance
    manifest = extract_transcript(mock_prs_instance, cue_token=DEFAULT_CUE)
    assert manifest["slides"][0]["segments"] == [
        {"kind": "text", "text": "This is [some text] not a cue. But this is"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "a cue."}
    ]

    # Scenario 11: Cue with surrounding whitespace in notes
    mock_prs_instance = mock_slide_with_notes([f"Before \t {DEFAULT_CUE} \n After"])
    mock_pptx_presentation.return_value = mock_prs_instance
    manifest = extract_transcript(mock_prs_instance, cue_token=DEFAULT_CUE)
    assert manifest["slides"][0]["segments"] == [
        {"kind": "text", "text": "Before"},
        {"kind": "cue", "cue": DEFAULT_CUE},
        {"kind": "text", "text": "After"}
    ]

    # Scenario 12: User reported issue - brackets in text, cue token is "[transition]"
    # Text: "Text before. [transition] Text after."
    # Cue Token: "[transition]"
    # Expected cue: "[transition]", Expected text segments: "Text before." and "Text after."
    notes_text_user_issue = "Text before. [transition] Text after."
    mock_prs_instance = mock_slide_with_notes([notes_text_user_issue])
    mock_pptx_presentation.return_value = mock_prs_instance
    manifest = extract_transcript(mock_prs_instance, cue_token=DEFAULT_CUE) # DEFAULT_CUE is "[transition]"
    assert manifest["slides"][0]["segments"] == [
        {"kind": "text", "text": "Text before."},
        {"kind": "cue", "cue": DEFAULT_CUE}, # Should be "[transition]"
        {"kind": "text", "text": "Text after."}
    ], f"Failed Scenario 12 for input: {notes_text_user_issue}"
